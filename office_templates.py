# Presentation imports
from pptx import Presentation
from pptx.util import Inches
from os import listdir, getlogin

# Email imports
import win32com.client

# Word Imports
import docx

global DEFAULT_SAVE_PATH

DEFAULT_SAVE_PATH = f"C:/Users/{getlogin()}/Desktop"


def create_image_list(images,path):
    """
    Create a list of image file paths by appending the path to each image file name.

    Args:
        images: list of str
        The list of image file names.
        path: str
        The path to the folder containing the images.

    Returns:
        list of str: The list of image file paths.
    """
    images_list = []

    for image in images:
        if image.endswith(('.jpg', '.jpeg', '.png', '.gif')):
            images_list.append(path+"/"+image)
        
    return images_list


def add_text(slide,text):
    """
    Add text to the slide placeholder.

    Args:
        slide: pptx.slide.Slide
        The slide object.
        text: str
        The text to be added to the slide.
    """

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        else:
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = text

        
def presentation(images, path="",title="Presentation",sizes=False,save_path=DEFAULT_SAVE_PATH):
    """
    Creates a PowerPoint presentation in SPX Template using the images provided.
    The slide order is defined by the images list (or folder) order.

    Args:
        images (str or list): The images to include in the presentation.
            - str: Path to the images folder.
            - list: Paths to each image file.

        path (str): The path to the images folder (only used if images is a list).

        title (str): The title of the presentation.

        sizes (list of tuples, optional): 
            A list of tuples where each tuple represents the (width, height) 
            for each image in the images list (in the same order).

        save_path (str): The path to save the PowerPoint presentation.

    Returns:
        None

    Raises:
        AssertionError: If the images parameter is not a list or a string.
    """

    assert type(images) == str or type(images) == list, "image should be a list (of images files names) or a string (the images folder path)"

    if type(images) == str:
        images = create_image_list(listdir(images),images)

    elif type(images) == list:
        if path != "": images = create_image_list(images,path)

    prs = Presentation(rf"template.pptx")
    
    #this is the order of the slides in the template (you might need to change it)
    inicial = prs.slide_layouts[0]
    content = prs.slide_layouts[3]
    final = prs.slide_layouts[4]

    slide = prs.slides.add_slide(inicial)
    add_text(slide,title)

    if sizes:

        for j, i in enumerate(images):
            slide = prs.slides.add_slide(content)

            if sizes[j] == 1:slide.shapes.add_picture(i,left=Inches(1), top=Inches(0.7), width=Inches(8), height=Inches(4.5))
            elif  sizes[j] == 2: slide.shapes.add_picture(i,left=Inches(2), top=Inches(0.7), width=Inches(6), height=Inches(4.5))
            elif  sizes[j] == 3: slide.shapes.add_picture(i,left=Inches(1), top=Inches(0.7), width=Inches(8), height=Inches(4.5))
            elif  sizes[j] == 4: slide.shapes.add_picture(i,left=Inches(0.5), top=Inches(0.7), width=Inches(9), height=Inches(4.5))
            elif  sizes[j] == 5: slide.shapes.add_picture(i,left=Inches(0), top=Inches(0.7),width=Inches(0.1), height=Inches(0.1))

    else:
        for i in images:
            slide = prs.slides.add_slide(content)
            slide.shapes.add_picture(i,left=Inches(2), top=Inches(0.7), width=Inches(8), height=Inches(4.5))


    prs.slides.add_slide(final)
    prs.save(f"{save_path}/{title.replace(' ','_')}.pptx")


def email(images=None,body='',path="",to_email="",subject="",attachments=False,send=True,images_on_body=True):
    """
    Creates an Outlook email with all images and subject provided and sends it to the emails in the to_email parameter.

    Args:
        images (str or list or None or bool): The images to include in the email.
            - str: Path to the images folder.
            - list: Paths to each image file.
            - None or False: If no image is needed.

        body (str): The body of the email.

        path (str): The path to the images folder (only used if images is a list).

        to_email (str): The email address to send the email to.
            - For multiple emails, separate them with ';', for example: "myemailaddress@gmail.com; otheremailaddress@outlook.com".

        subject (str): The email subject.

        attachments (str or list or bool, optional): Path for the attachments.
            - Use string if it's only one attachment.
            - Use list for multiple attachments.
            - Use False if there are no attachments. Default is False.

        send (bool, optional): Whether to send the email.
            - If True, sends the email.
            - If False, opens the email as a draft. Note that the function only stops running when the Outlook window is closed. Default is True.

    Returns:
        None

    Raises:
        AssertionError: If the images parameter is not a list or a string.
    """

    outlook = win32com.client.Dispatch('outlook.application')
    mail = outlook.CreateItem(0)
    mail.Subject = subject
    mail.To = to_email
    
    if images:

        assert type(images) == str or type(images) == list, "image should be a list (of images files names) or a string (the images folder path)"
        
        if type(images) == str:
            images = create_image_list(listdir(images),images)
            
        elif type(images) == list:
            if path != "": images = create_image_list(images,path)
            

        image_template = '<img src="cid:#" width="900" height="500" /><br>'

        counter = 0

        for image in images: 
            attachment = mail.Attachments.Add(image)
            img_id = f"Img{counter}"
            attachment.PropertyAccessor.SetProperty("http://schemas.microsoft.com/mapi/proptag/0x3712001F", img_id)
            if images_on_body: body +=  image_template.replace("#",img_id) 
            counter += 1
            
    mail.HTMLBody  = body

    if attachments:
        if type(attachments) == str:
            attachment = mail.Attachments.Add(attachments)
            
        elif type(attachments) == list:
            for at in attachments:
                attachment = mail.Attachments.Add(at)

    if send:
        mail.Send()
    else:
        mail.Display(True)
        

def word(images, path="",title="Presentation",sizes=False,save_path=DEFAULT_SAVE_PATH):
    """
    Creates a Word document using the images provided.
    The order of the images in the document is defined by the order in the images list.

    Args:
        images (str or list): The images to include in the document.
            - str: Path to the images folder.
            - list: Paths to each image file.
        path (str): The path to the images folder (only used if images is a list).
        title (str): The title of the Word document.
        sizes (list of tuples): A list of tuples where each tuple represents the (width, height) for each image in the images list (in the same order).
        save_path (str): The path to save the Word document.

    Returns:
        None

    Raises:
        AssertionError: If the images parameter is not a list or a string.
    """

    assert type(images) == str or type(images) == list, "image should be a list (of images files names) or a string (the images folder path)"
    
    if type(images) == str:
        images = create_image_list(listdir(images),images)

    elif type(images) == list:
        if path != "": images = create_image_list(images,path)

    my_doc = docx.Document()

    if sizes:
        for j, i in enumerate(images):
            if sizes[j] == 1: my_doc.add_picture(i, width=docx.shared.Inches(5), height=docx.shared.Inches(3))
            elif sizes[j] == 2: my_doc.add_picture(i, width=docx.shared.Inches(4), height=docx.shared.Inches(4))
            elif sizes[j] == 3: my_doc.add_picture(i, width=docx.shared.Inches(5), height=docx.shared.Inches(4))

    else:
        for i in images:
            try:
                my_doc.add_picture(i, width=docx.shared.Inches(4), height=docx.shared.Inches(4))
            except: continue

    my_doc.save(f"{save_path}/{title}.docx")
